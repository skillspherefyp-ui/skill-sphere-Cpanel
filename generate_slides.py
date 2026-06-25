import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# Theme Colors
NAVY        = RGBColor(0x0D, 0x0D, 0x1E)
NAVY2       = RGBColor(0x1A, 0x1A, 0x2E)
NAVY3       = RGBColor(0x25, 0x25, 0x40)
ORANGE      = RGBColor(0xFF, 0x8C, 0x42)
WHITE       = RGBColor(0xF8, 0xFA, 0xFC)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
ACCENT_BLUE = RGBColor(0x4F, 0x46, 0xE5)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
SLATE       = RGBColor(0x64, 0x74, 0x8B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_tb(slide, text, x, y, w, h,
           size=16, bold=False, color=None, align=PP_ALIGN.LEFT,
           italic=False, name="Segoe UI"):
    if color is None:
        color = WHITE
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = name
    return tb

def add_tb_lines(slide, lines, x, y, w, h,
                 default_size=14, default_color=None, name="Segoe UI"):
    if default_color is None:
        default_color = WHITE
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, bold, sz, col = item, False, default_size, default_color
        else:
            txt  = item[0]
            bold = item[1] if len(item) > 1 else False
            sz   = item[2] if len(item) > 2 else default_size
            col  = item[3] if len(item) > 3 else default_color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = col
        r.font.name = name
    return tb

def slide_bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)

def header_band(slide, title):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), NAVY2)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.04), ORANGE)
    add_tb(slide, title,
           Inches(0.5), Inches(0.2), Inches(12), Inches(0.75),
           size=28, bold=True, color=ORANGE)

def footer_bar(slide):
    add_rect(slide, 0, SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), NAVY2)
    add_tb(slide, "SkillSphere  An AI-Powered Personalized Education and Learning Platform",
           Inches(0.4), SLIDE_H - Inches(0.38), Inches(12), Inches(0.35),
           size=10, color=MUTED)

def slide_num(slide, n):
    add_tb(slide, str(n),
           SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.44),
           Inches(0.6), Inches(0.35),
           size=12, color=MUTED, align=PP_ALIGN.RIGHT)

def set_cell_bg(cell, rgb):
    r, g, b = rgb[0], rgb[1], rgb[2]
    hex_color = f"{r:02X}{g:02X}{b:02X}"
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', hex_color)

def set_cell_text(cell, text, size=12, bold=False,
                  color=None, align=PP_ALIGN.LEFT, name="Segoe UI"):
    if color is None:
        color = WHITE
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = name

def set_cell_border_none(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for side in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = etree.SubElement(tcPr, qn(side))
        ln.set('w', '0')
        noFill = etree.SubElement(ln, qn('a:noFill'))

def add_req_table(slide, rows_data, x, y, w, h):
    """
    rows_data: list of (sno, requirement, status)
    Header row + data rows
    """
    total_rows = len(rows_data) + 1  # +1 for header
    table = slide.shapes.add_table(total_rows, 3, x, y, w, h).table

    # Column widths
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(8.5)
    table.columns[2].width = Inches(1.6)

    # Header row
    headers = ["S.No", "Functional Requirement", "Status"]
    header_bg = (0x0D, 0x0D, 0x1E)  # NAVY as tuple for set_cell_bg
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        set_cell_bg(cell, (0xFF, 0x8C, 0x42))  # ORANGE header
        set_cell_text(cell, hdr, size=13, bold=True,
                      color=RGBColor(0x0D, 0x0D, 0x1E),
                      align=PP_ALIGN.CENTER)

    # Data rows
    for ri, (sno, req, status) in enumerate(rows_data, start=1):
        row_bg = (0x1A, 0x1A, 0x2E) if ri % 2 == 1 else (0x25, 0x25, 0x40)

        cell_sno = table.cell(ri, 0)
        set_cell_bg(cell_sno, row_bg)
        set_cell_text(cell_sno, str(sno), size=12, bold=True,
                      color=ORANGE, align=PP_ALIGN.CENTER)

        cell_req = table.cell(ri, 1)
        set_cell_bg(cell_req, row_bg)
        set_cell_text(cell_req, req, size=11.5, color=WHITE)

        cell_st = table.cell(ri, 2)
        set_cell_bg(cell_st, row_bg)
        st_color = GREEN if status == "Completed" else AMBER
        set_cell_text(cell_st, status, size=11.5, bold=True,
                      color=st_color, align=PP_ALIGN.CENTER)

# ── 31 Functional Requirements (Super Admin→Admin, Admin→Instructor) ──────────
ALL_REQS = [
    (1,  "Admin creates Instructors & Experts",                                                              "Completed"),
    (2,  "Admin Activates/deactivates Instructor & Expert accounts",                                         "Completed"),
    (3,  "Instructor and Expert Logins using verified Credentials",                                          "Completed"),
    (4,  "Instructor Creates Skill categories.",                                                              "Completed"),
    (5,  "Instructor Creates Skill courses, adds description, level, and language, selects course delivery mode: Manual (PDF/topic-based) or AI-Generated, and uses prompt engineering buttons to customize AI lecture generation tone, depth, and style.", "Completed"),
    (6,  "Instructor adds topics and structured outlines.",                                                   "Completed"),
    (7,  "Instructor uploads course materials",                                                              "Completed"),
    (8,  "Instructor submits course to AI for lecture generation.",                                          "Completed"),
    (9,  "Instructor publish/unpublish courses",                                                             "Completed"),
    (10, "Instructor reviews and publishes the AI-generated course.",                                        "Completed"),
    (11, "Instructor updates or deletes an existing course.",                                                "Completed"),
    (12, "Expert review and provide Feedback of Courses",                                                    "Completed"),
    (13, "Student signs up using email, password, and personal details",                                     "Completed"),
    (14, "Student verifies email using OTP.",                                                                "Completed"),
    (15, "Student logs in using verified email and password.",                                               "Completed"),
    (16, "Admin Activates/Deactivates Students",                                                             "Completed"),
    (17, "Student browses available skill categories.",                                                      "Completed"),
    (18, "Student selects a skill and views all published courses under it.",                                "Completed"),
    (19, "Student views the course outline (locked/unlocked topics).",                                       "Completed"),
    (20, "Student starts the first unlocked topic.",                                                         "Completed"),
    (21, "Student interacts with the AI tutor using voice/text-based questions.",                            "Completed"),
    (22, "AI pauses the lecture, answers the question, and resumes automatically.",                          "Completed"),
    (23, "Student attempts AI-generated quizzes (MCQs).",                                                    "Completed"),
    (24, "Student views quiz results and AI-generated feedback.",                                            "Completed"),
    (25, "Student downloads AI-generated flashcards after payment.",                                         "Completed"),
    (26, "Student pays for the certificate through a secure payment method.",                                "Pending"),
    (27, "Student views overall course progress in the dashboard.",                                          "Completed"),
    (28, "Student receives the certificate in downloadable PDF format.",                                     "Completed"),
    (29, "Student enters AI Assistant mode to ask additional open-ended questions.",                         "Completed"),
    (30, "Admins, Instructors, Experts and Students Receive Emails.",                                        "Completed"),
]

# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
import io
from PIL import Image as PILImage, ImageDraw

# --- Build real PNG placeholder so right-click → "Change Picture" ---
IMG_PX_W, IMG_PX_H = 830, 1500
img = PILImage.new("RGB", (IMG_PX_W, IMG_PX_H), (0x1A, 0x1A, 0x2E))
draw = ImageDraw.Draw(img)
brd = 16
# Orange border frame
draw.rectangle([0, 0, IMG_PX_W-1, IMG_PX_H-1],
               outline=(0xFF, 0x8C, 0x42), width=brd)
# Subtle diagonal cross
lc = (0x25, 0x25, 0x50)
draw.line([(brd, brd), (IMG_PX_W-brd, IMG_PX_H-brd)], fill=lc, width=2)
draw.line([(IMG_PX_W-brd, brd), (brd, IMG_PX_H-brd)], fill=lc, width=2)
# Centre icon: circle + triangle (landscape metaphor)
cx, cy = IMG_PX_W // 2, IMG_PX_H // 2
draw.ellipse([cx-75, cy-240, cx+75, cy-90],
             outline=(0xFF, 0x8C, 0x42), width=5)
draw.polygon([(cx-210, cy+100), (cx, cy-60), (cx+210, cy+100)],
             outline=(0x94, 0xA3, 0xB8), width=4)
draw.polygon([(cx-80, cy+100), (cx+130, cy+10), (cx+310, cy+100)],
             outline=(0x64, 0x74, 0x8B), width=4)
# Bottom label strip
draw.rectangle([brd, IMG_PX_H-100, IMG_PX_W-brd, IMG_PX_H-brd],
               fill=(0x0D, 0x0D, 0x1E))
img_bytes = io.BytesIO()
img.save(img_bytes, format="PNG")
img_bytes.seek(0)

# ─── Slide setup ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)  # full navy

# Vertical measurements — nothing overlaps
# Slide height = 7.5"
# Left content width = 8.75" (to x=8.88"), right image = 4.45"

CONTENT_X  = Inches(0.38)   # left text margin (after strip)
CONTENT_W  = Inches(8.15)   # text width

# ── Structural background shapes (no overlap with any text) ──────────
# Thin orange strip on very left edge
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, ORANGE)

# Top orange bar
add_rect(s, 0, 0, SLIDE_W, Inches(0.07), ORANGE)

# Bottom card: y = 4.70" → 7.50"  (height = 2.80")
CARD_Y = Inches(4.70)
add_rect(s, 0, CARD_Y, Inches(8.88), SLIDE_H - CARD_Y, NAVY2)
add_rect(s, 0, CARD_Y, Inches(8.88), Inches(0.05), ORANGE)  # card top border

# ── Right 1/3 — real picture placeholder ────────────────────────────
IMG_X = Inches(8.88)
s.shapes.add_picture(img_bytes, IMG_X, Inches(0), SLIDE_W - IMG_X, SLIDE_H)

# ── Left 2/3 — text, carefully spaced, no element crosses another ────

# Row 1  y=0.45"   SKILLSPHERE  60pt  height≈1.3"  → ends ≈1.75"
add_tb(s, "SKILLSPHERE",
       CONTENT_X, Inches(0.45), CONTENT_W, Inches(1.35),
       size=60, bold=True, color=ORANGE)

# Row 2  y=1.80"   Subtitle  19pt  height≈0.5"  → ends ≈2.30"
add_tb(s, "An AI-Powered Personalized Education and Learning Platform",
       CONTENT_X, Inches(1.80), CONTENT_W, Inches(0.55),
       size=19, color=WHITE)

# Row 3  y=2.52"   Orange rule (starts after subtitle ends + gap)
add_rect(s, CONTENT_X, Inches(2.52), Inches(5.2), Inches(0.05), ORANGE)

# Row 4  y=2.72"   Presented by  → ends ≈3.35"
add_tb_lines(s,
    [("Presented by:", True, 12, MUTED),
     ("Talha Rizwan   |   Muhammad Danish   |   Tehzeen Abbas", False, 15, WHITE)],
    CONTENT_X, Inches(2.72), CONTENT_W, Inches(0.70))

# Row 5  y=3.55"   Supervised by  → ends ≈4.20"
add_tb_lines(s,
    [("Supervised by:", True, 12, MUTED),
     ("Sir Salman Ahmed", False, 15, WHITE)],
    CONTENT_X, Inches(3.55), CONTENT_W, Inches(0.70))

# Row 6  y=4.40"   Second rule (0.25" gap before card starts at 4.70")
add_rect(s, CONTENT_X, Inches(4.40), Inches(3.5), Inches(0.05), ORANGE)

# ── Inside bottom card (y > 4.75") ──────────────────────────────────

# Row 7  y=4.85"   Department of Computer Science  14pt  → ends ≈5.15"
add_tb(s, "Department of Computer Science",
       CONTENT_X, Inches(4.85), CONTENT_W, Inches(0.38),
       size=14, bold=True, color=WHITE)

# Row 8  y=5.25"   University  12pt  → ends ≈5.55"
add_tb(s, "Capital University of Science and Technology, Islamabad",
       CONTENT_X, Inches(5.25), CONTENT_W, Inches(0.38),
       size=12, color=MUTED)

# Row 9  y=5.75"   Thin rule inside card
add_rect(s, CONTENT_X, Inches(5.75), Inches(4.5), Inches(0.03), NAVY3)

# Row 10  y=5.95"  Project label  11pt
add_tb(s, "Final Year Project  |  2026",
       CONTENT_X, Inches(5.98), CONTENT_W, Inches(0.35),
       size=11, color=SLATE)

# Row 11  y=6.50"  Decorative dots
add_rect(s, CONTENT_X,                Inches(6.55), Inches(0.22), Inches(0.22), ORANGE)
add_rect(s, CONTENT_X + Inches(0.35), Inches(6.55), Inches(0.22), Inches(0.22), NAVY3)
add_rect(s, CONTENT_X + Inches(0.70), Inches(6.55), Inches(0.22), Inches(0.22), NAVY3)

# Row 12  y=7.05"  Bottom tag line
add_tb(s, "An LMS powered by AI",
       CONTENT_X, Inches(7.05), CONTENT_W, Inches(0.35),
       size=11, italic=True, color=SLATE)

# ── SLIDE 2 — Functional Requirements (1–16) ─────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_band(s, "Functional Requirements")
footer_bar(s)
slide_num(s, 2)

add_req_table(
    s, ALL_REQS[:15],
    x=Inches(0.4), y=Inches(1.2),
    w=Inches(12.5), h=Inches(5.9)
)

# ── SLIDE 3 — Functional Requirements (17–31) ────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_band(s, "Functional Requirements (cont.)")
footer_bar(s)
slide_num(s, 3)

add_req_table(
    s, ALL_REQS[15:],
    x=Inches(0.4), y=Inches(1.2),
    w=Inches(12.5), h=Inches(5.5)
)

# ── SLIDE 4 — User Roles ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_band(s, "User Roles")
footer_bar(s)
slide_num(s, 4)

roles = [
    ("Student",     ORANGE,
     "Register & login with OTP verification\nEnroll in courses\nLearn Course\nAsk 24/7 AI Assistant\nTrack progress & earn certificates\nParticipate in discussion threads\nManage personal to-do tasks"),
    ("Expert",      ACCENT_BLUE,
     "Receive credentials via email from Admin\nReview course outlines & topics\nValidate uploaded PDF materials\nProvide structured quality feedback\nEnsure content accuracy before publishing"),
    ("Instructor",  GREEN,
     "Receive credentials via email from Admin\nCreate & manage courses\nUpload outlines, topics & PDF materials\nMonitor student progress\nReceive expert feedback on courses"),
    ("Admin",       AMBER,
     "Admin has Full Access of Instructor with some additional responsibilities\nCreate Instructor & Expert accounts\nSend credentials via email\nManage all users & permissions\nManage course categories\nOversee certificates & platform analytics"),
]

col_w  = Inches(2.9)
col_h  = Inches(4.9)
gap    = Inches(0.25)
sx     = Inches(0.4)
sy     = Inches(1.35)

for i, (role, accent, desc) in enumerate(roles):
    cx = sx + i * (col_w + gap)
    add_rect(s, cx, sy, col_w, col_h, NAVY2)
    add_rect(s, cx, sy, col_w, Inches(0.08), accent)
    add_tb(s, role, cx + Inches(0.18), sy + Inches(0.18),
           col_w - Inches(0.3), Inches(0.5),
           size=20, bold=True, color=accent)
    add_tb(s, desc, cx + Inches(0.18), sy + Inches(0.82),
           col_w - Inches(0.3), col_h - Inches(1.0),
           size=12.5, color=WHITE)

# ── SLIDE 5 — Workflow ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_band(s, "Platform Workflow")
footer_bar(s)
slide_num(s, 5)

steps = [
    ("1", "Admin\nCreates Accounts",
     "Admin creates Instructor & Expert accounts via dashboard. Credentials are sent automatically by email."),
    ("2", "Instructor\nCreates Courses",
     "Instructor uploads outlines & PDFs, adds topics, and chooses delivery mode: Manual (PDF/topic-based) or AI-Generated."),
    ("3", "Expert\nReviews Courses",
     "Expert reviews course outlines, validates uploaded materials and provides structured feedback."),
    ("4", "Student\nEnrolls & Learns",
     "Student enrolls and accesses the course in the mode already set by the Instructor — Manual or AI."),
    ("5", "AI or Manual\nContent Delivery",
     "Manual Mode: student reads Instructor-uploaded PDFs & notes. AI Mode: AI generates lectures, audio & diagrams."),
    ("6", "Quizzes &\nAssignments",
     "Auto-generated quizzes and assignments after each topic. Students can interact with AI anytime."),
    ("7", "Certificate\nIssued",
     "On 100% course completion, a personalized PDF certificate is auto-generated and emailed."),
    ("8", "Email\nNotifications",
     "Welcome emails, OTP verification, credential delivery, and certificate emails — all automated."),
]

box_w  = Inches(2.9)
box_h  = Inches(2.55)
gap    = Inches(0.22)
row1_y = Inches(1.26)
row2_y = row1_y + box_h + Inches(0.24)
sx     = Inches(0.28)

for i, (num, title, desc) in enumerate(steps):
    row = i // 4
    col = i % 4
    bx  = sx + col * (box_w + gap)
    by  = row1_y if row == 0 else row2_y

    add_rect(s, bx, by, box_w, box_h, NAVY2)
    add_rect(s, bx, by, box_w, Inches(0.07), ORANGE)
    add_tb(s, num, bx + Inches(0.15), by + Inches(0.1),
           Inches(0.4), Inches(0.38), size=16, bold=True, color=ORANGE)
    add_tb(s, title, bx + Inches(0.15), by + Inches(0.44),
           box_w - Inches(0.3), Inches(0.7), size=13, bold=True, color=WHITE)
    add_tb(s, desc, bx + Inches(0.15), by + Inches(1.14),
           box_w - Inches(0.3), Inches(1.28), size=10.5, color=MUTED)

add_tb(s, "  \u25bc  workflow continues  \u25bc",
       Inches(5.6), row1_y + box_h + Inches(0.02),
       Inches(2.5), Inches(0.22), size=10, color=ORANGE, align=PP_ALIGN.CENTER)

# ── SLIDE 6 — Courses Data Structure ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_band(s, "Courses Data Structure")
footer_bar(s)
slide_num(s, 6)

data_sections = [
    ("Categories", ORANGE, [
        "Name",
        "Linked to multiple Courses",
    ]),
    ("Courses", ACCENT_BLUE, [
        "Name & Description",
        "Thumbnail / Cover Image",
        "Language & Duration",
        "Level: Beginner / Intermediate / Advanced",
        "Category (linked to Category Name)",
        "Topics (multiple per course)",
        "Materials: PDFs (outlines, notes, resources)",
    ]),
    ("Topics", GREEN, [
        "Name",
        "Linked to Course",
        "Materials: PDFs (notes, exercises, references)",
    ]),
    ("Materials", AMBER, [
        "Available for both Courses and Topics",
        "PDF format only",
        "Outlines, lecture notes",
        "Exercises & reference documents",
    ]),
]

col_w  = Inches(2.9)
col_h  = Inches(4.95)
gap    = Inches(0.25)
sx     = Inches(0.4)
sy     = Inches(1.3)

for i, (title, accent, items) in enumerate(data_sections):
    cx = sx + i * (col_w + gap)
    add_rect(s, cx, sy, col_w, col_h, NAVY2)
    add_rect(s, cx, sy, col_w, Inches(0.08), accent)
    add_tb(s, title, cx + Inches(0.18), sy + Inches(0.15),
           col_w - Inches(0.3), Inches(0.45), size=17, bold=True, color=accent)
    add_rect(s, cx + Inches(0.18), sy + Inches(0.66),
             col_w - Inches(0.36), Inches(0.02), accent)
    bullets_text = "\n".join("\u2022  " + item for item in items)
    add_tb(s, bullets_text,
           cx + Inches(0.18), sy + Inches(0.82),
           col_w - Inches(0.3), col_h - Inches(0.95),
           size=12, color=WHITE)

# ── SLIDE 7 — SEO Strategy ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_band(s, "Our Multi-Faceted SEO Approach")
footer_bar(s)
slide_num(s, 7)

add_tb(s,
       "SkillSphere employs a blended strategy covering four key pillars to maximize search engine performance.",
       Inches(0.5), Inches(1.18), Inches(12.3), Inches(0.45),
       size=13, color=MUTED, italic=True)

pillars = [
    ("Technical SEO", ORANGE, [
        "Static content injection: hidden H1/H2/UL in index.html — Google indexes course catalog even in React CSR",
        "sitemap.xml on cPanel — 6 URLs: / , /explore , /blog , /about , /help , /certificate-verify with priority & changefreq",
        "robots.txt on cPanel — Allow: * (all pages crawlable), points to sitemap.xml",
        "Full meta tag set: description, keywords (100+ Pakistan terms), robots, author",
        "Open Graph tags: og:title, og:description, og:image, og:site_name",
        "Twitter Card: summary_large_image with title, description & OG image",
        "font-display:swap on all custom fonts + Google Fonts preconnect",
        "Deep linking — /explore & course routes are directly shareable URLs",
    ]),
    ("On-Page SEO", ACCENT_BLUE, [
        "Dynamic document.title updated per page via JavaScript",
        "setMeta() injects OG & Twitter tags at runtime per screen",
        "og-image.png — custom social preview image for all shares",
        "Canonical URLs on: / , /explore , /blog , /blog/:id , /help",
        "Keywords meta: 100+ Pakistan-specific terms targeting learner searches",
    ]),
    ("Content SEO", GREEN, [
        "Blog system — BlogScreen + BlogPostScreen with per-post canonical URL",
        "Help Center page (/help) — FAQ-style long-tail keyword coverage",
        "Static course catalog in HTML: 100+ course names across 8 categories",
        "About, Community, Privacy Policy & Terms pages for trust signals",
        "All pages in English with Urdu learner-targeting keyword set",
    ]),
    ("Off-Page SEO", AMBER, [
        "Share modal on course detail — share to WhatsApp, Facebook, Twitter",
        "OG image ensures rich social previews on all platforms",
        "Email system: welcome, OTP, credentials & certificates drive return visits",
        "Certificates delivered as PDF — shareable on LinkedIn & resumes",
    ]),
]

col_w  = Inches(2.9)
col_h  = Inches(5.3)
gap    = Inches(0.25)
sx     = Inches(0.4)
sy     = Inches(1.55)

for i, (title, accent, items) in enumerate(pillars):
    cx = sx + i * (col_w + gap)
    add_rect(s, cx, sy, col_w, col_h, NAVY2)
    add_rect(s, cx, sy, col_w, Inches(0.08), accent)
    add_tb(s, title, cx + Inches(0.18), sy + Inches(0.15),
           col_w - Inches(0.3), Inches(0.45), size=16, bold=True, color=accent)
    add_rect(s, cx + Inches(0.18), sy + Inches(0.66),
             col_w - Inches(0.36), Inches(0.02), accent)
    bullets_text = "\n".join("\u2022  " + item for item in items)
    add_tb(s, bullets_text,
           cx + Inches(0.18), sy + Inches(0.82),
           col_w - Inches(0.3), col_h - Inches(0.95),
           size=10, color=WHITE)

# ── SLIDE 8 — Thank You ───────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s)

add_rect(s, SLIDE_W - Inches(4), 0, Inches(4), Inches(3), NAVY2)
add_rect(s, SLIDE_W - Inches(2.5), 0, Inches(2.5), Inches(2), NAVY3)
add_rect(s, 0, SLIDE_H - Inches(2), Inches(4.5), Inches(2), NAVY2)
add_rect(s, 0, SLIDE_H - Inches(2), Inches(0.12), Inches(2), ORANGE)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, ORANGE)

add_tb(s, "THANK YOU",
       Inches(0.5), Inches(1.9), Inches(10), Inches(2.0),
       size=72, bold=True, color=ORANGE)

add_rect(s, Inches(0.5), Inches(3.82), Inches(7), Inches(0.06), ORANGE)

add_tb(s, "that brings us to the end.",
       Inches(0.5), Inches(4.05), Inches(10), Inches(0.5),
       size=18, italic=True, color=WHITE)

add_tb(s, "We'd like to thank you for your time, interest and attention today.",
       Inches(0.5), Inches(4.58), Inches(10), Inches(0.5),
       size=15, color=MUTED)

add_tb_lines(s,
    [("Talha Rizwan  |  Muhammad Danish  |  Tehzeen Abbas", False, 13, MUTED),
     ("Supervised by Sir Salman Ahmed", False, 12, SLATE)],
    Inches(0.5), SLIDE_H - Inches(1.6), Inches(9), Inches(0.9))

# ── Save ──────────────────────────────────────────────────────────────────────
out = "D:/skillsphere design Project part 2/SkillSphere_Presentation_Updated.pptx"
prs.save(out)

# ── Embed thumbnail so Windows Explorer shows a preview ───────────────────────
import zipfile, os

thumb_w, thumb_h = 640, 360
thumb_img = PILImage.new("RGB", (thumb_w, thumb_h), (0x0D, 0x0D, 0x1E))
td = ImageDraw.Draw(thumb_img)
td.rectangle([0, 0, 18, thumb_h],      fill=(0xFF, 0x8C, 0x42))  # left strip
td.rectangle([0, 0, thumb_w, 8],       fill=(0xFF, 0x8C, 0x42))  # top bar
td.rectangle([thumb_w-220, 0, thumb_w, thumb_h], fill=(0x1A, 0x1A, 0x2E))  # right card
try:
    fnt_big   = ImageFont.truetype("arialbd.ttf", 54)
    fnt_small = ImageFont.truetype("arial.ttf",   20)
    td.text((38, 75),  "SKILLSPHERE",                    fill=(0xFF,0x8C,0x42), font=fnt_big)
    td.text((38, 150), "An AI-Powered Personalized",     fill=(0xF8,0xFA,0xFC), font=fnt_small)
    td.text((38, 175), "Education and Learning Platform",fill=(0xF8,0xFA,0xFC), font=fnt_small)
except Exception:
    pass
td.rectangle([38, 208, 320, 213], fill=(0xFF, 0x8C, 0x42))  # underline

thumb_bytes_val = io.BytesIO()
thumb_img.save(thumb_bytes_val, "JPEG", quality=85)
thumb_jpeg = thumb_bytes_val.getvalue()

# Read original zip, patch Content_Types, inject thumbnail — all in one pass
ct_name   = '[Content_Types].xml'
thumb_ct  = '<Default Extension="jpeg" ContentType="image/jpeg"/>'
thumb_rel = '<Relationship Id="rId1000" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/thumbnail" Target="docProps/thumbnail.jpeg"/>'

with zipfile.ZipFile(out, 'r') as zin:
    entries = {item.filename: zin.read(item.filename) for item in zin.infolist()}

# Patch Content_Types.xml
ct_xml = entries[ct_name].decode('utf-8')
if thumb_ct not in ct_xml:
    ct_xml = ct_xml.replace('</Types>', thumb_ct + '</Types>')
entries[ct_name] = ct_xml.encode('utf-8')

# Patch _rels/.rels to add thumbnail relationship
rels_name = '_rels/.rels'
if rels_name in entries:
    rels_xml = entries[rels_name].decode('utf-8')
    if 'thumbnail' not in rels_xml:
        rels_xml = rels_xml.replace('</Relationships>', thumb_rel + '</Relationships>')
    entries[rels_name] = rels_xml.encode('utf-8')

# Add thumbnail image (replace old if exists)
entries['docProps/thumbnail.jpeg'] = thumb_jpeg

# Write patched zip
tmp = out + ".tmp"
with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
    for name, data in entries.items():
        zout.writestr(name, data)
os.replace(tmp, out)

print("Saved:", out)
print(f"Total slides: {len(prs.slides)}")
print("Thumbnail embedded — Windows Explorer preview ready.")
